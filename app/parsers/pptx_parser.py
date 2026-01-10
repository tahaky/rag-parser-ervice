from typing import Dict, Any, List
from pptx import Presentation

from app.parsers.base import BaseParser


class PptxParser(BaseParser):
    """Parser for PPTX presentations."""

    def parse(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX presentation into structured JSON."""
        prs = Presentation(file_path)

        # Extract metadata
        core_props = prs.core_properties
        metadata = self.build_metadata(
            title=core_props.title or "",
            author=core_props.author or "",
            subject=core_props.subject or "",
            created=str(core_props.created) if core_props.created else None,
            modified=str(core_props.modified) if core_props.modified else None,
        )

        # Parse slides
        slides = []
        total_text_length = 0
        total_tables = 0
        total_images = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "tables": [],
                "notes": "",
            }

            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text.strip()

            # Extract content from shapes
            for shape in slide.shapes:
                # Text content
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and text != slide_data["title"]:
                        content_item = {
                            "type": "text",
                            "text": text,
                            "hash": self.calculate_hash(text),
                        }
                        slide_data["content"].append(content_item)
                        total_text_length += len(text)

                # Tables
                if shape.has_table:
                    table_data = self._extract_table(shape.table)
                    slide_data["tables"].append(table_data)
                    total_tables += 1

                # Images
                if hasattr(shape, "image"):
                    total_images += 1

            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["notes"] = notes_text
                    slide_data["notes_hash"] = self.calculate_hash(notes_text)
                    total_text_length += len(notes_text)

            slides.append(slide_data)

        # Build structure
        structure = {
            "format": "pptx",
            "slides": slides,
        }

        # Build stats
        stats = self.build_stats(
            total_pages=len(slides),
            total_text_length=total_text_length,
            total_tables=total_tables,
            total_images=total_images,
            total_slides=len(slides),
        )

        return {
            "format": "pptx",
            "metadata": metadata,
            "structure": structure,
            "stats": stats,
        }

    def _extract_table(self, table) -> Dict[str, Any]:
        """Extract table data from PPTX."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        table_text = " ".join(" ".join(row) for row in rows)
        return {
            "rows": rows,
            "row_count": len(rows),
            "col_count": len(rows[0]) if rows else 0,
            "hash": self.calculate_hash(table_text),
        }
